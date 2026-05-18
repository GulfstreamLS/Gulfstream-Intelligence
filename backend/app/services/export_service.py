import io
from typing import Any

from docx import Document as DocxDocument
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGB
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    HRFlowable,
    KeepTogether,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.services.storage_service import storage_service

def _sanitize(text: str) -> str:
    """Replace problematic Unicode characters that ReportLab renders as black boxes."""
    replacements = {
        "­": "-",   # soft hyphen
        "‑": "-",   # non-breaking hyphen
        "‐": "-",   # hyphen
        "‒": "-",   # figure dash
        "–": "-",   # en dash
        "—": "-",   # em dash
        "―": "-",   # horizontal bar
        "‘": "'",   # left single quotation
        "’": "'",   # right single quotation
        "“": '"',   # left double quotation
        "”": '"',   # right double quotation
        "•": "*",   # bullet
        "…": "...", # ellipsis
        " ": " ",   # non-breaking space
    }
    for char, replacement in replacements.items():
        text = text.replace(char, replacement)
    return text


# ── PDF colour palette (light theme) ──────────────────────────────────────────
_P_INDIGO       = colors.HexColor("#6366F1")
_P_INDIGO_DARK  = colors.HexColor("#4338CA")
_P_INDIGO_LIGHT = colors.HexColor("#EEF2FF")
_P_SLATE_800    = colors.HexColor("#1E293B")
_P_SLATE_600    = colors.HexColor("#475569")
_P_SLATE_400    = colors.HexColor("#94A3B8")
_P_SLATE_50     = colors.HexColor("#F8FAFC")
_P_BORDER       = colors.HexColor("#E2E8F0")
_P_WHITE        = colors.white

_P_SEV = {
    "Critical": colors.HexColor("#EF4444"),
    "High":     colors.HexColor("#F97316"),
    "Medium":   colors.HexColor("#D97706"),
    "Low":      colors.HexColor("#059669"),
}
_P_SEV_BG = {
    "Critical": colors.HexColor("#FEF2F2"),
    "High":     colors.HexColor("#FFF7ED"),
    "Medium":   colors.HexColor("#FFFBEB"),
    "Low":      colors.HexColor("#F0FDF4"),
}

# ── PPT colour palette (light theme) ──────────────────────────────────────────
_INDIGO       = RGBColor(0x63, 0x66, 0xF1)   # #6366f1
_INDIGO_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)   # #eef2ff
_INDIGO_DARK  = RGBColor(0x43, 0x38, 0xCA)   # #4338ca
_SLATE_800    = RGBColor(0x1E, 0x29, 0x3B)
_SLATE_600    = RGBColor(0x47, 0x55, 0x69)
_SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)
_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
_BG           = RGBColor(0xFA, 0xFB, 0xFF)   # near-white slide background
_RED          = RGBColor(0xEF, 0x44, 0x44)
_RED_LIGHT    = RGBColor(0xFE, 0xF2, 0xF2)
_ORANGE       = RGBColor(0xF9, 0x73, 0x16)
_ORANGE_LIGHT = RGBColor(0xFF, 0xF7, 0xED)
_YELLOW       = RGBColor(0xD9, 0x77, 0x06)
_YELLOW_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
_GREEN        = RGBColor(0x05, 0x96, 0x69)
_GREEN_LIGHT  = RGBColor(0xF0, 0xFD, 0xF4)

_SEV_COLOR    = {"Critical": _RED,    "High": _ORANGE,    "Medium": _YELLOW,    "Low": _GREEN}
_SEV_BG_COLOR = {"Critical": _RED_LIGHT, "High": _ORANGE_LIGHT, "Medium": _YELLOW_LIGHT, "Low": _GREEN_LIGHT}
_PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
_PURPLE_DARK  = RGBColor(0x4C, 0x1D, 0x95)
_PURPLE_LIGHT = RGBColor(0xFA, 0xF5, 0xFF)
_RED_DARK     = RGBColor(0x7F, 0x1D, 0x1D)
_INDIGO_MID   = RGBColor(0x81, 0x8C, 0xF8)   # lighter indigo for decorative use
_GREEN_DARK2  = RGBColor(0x05, 0x50, 0x36)   # deep green for action headers

# Slide: 13.33" × 7.5"
_SW = Inches(13.33)
_SH = Inches(7.5)


# ── PPT helpers ────────────────────────────────────────────────────────────────

def _solid_bg(slide: Any, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_rect(slide: Any, left: float, top: float, w: float, h: float,
              fill: RGBColor, line: RGBColor | None = None) -> Any:
    shape = slide.shapes.add_shape(1, int(left), int(top), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide: Any, text: str,
              left: float, top: float, w: float, h: float,
              size: int, bold: bool = False,
              color: RGBColor = _SLATE_800,
              align: PP_ALIGN = PP_ALIGN.LEFT) -> Any:
    txb = slide.shapes.add_textbox(int(left), int(top), int(w), int(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _blank(prs: Presentation) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_footer(slide: Any, label: str = "") -> None:
    """Thin branded footer bar at the bottom of every content slide."""
    _add_rect(slide, 0, _SH - Inches(0.3), _SW, Inches(0.3), _SLATE_800)
    _add_text(slide, "Gulfstream Intelligence  ·  Regulatory Analysis",
              Inches(0.3), _SH - Inches(0.28), Inches(7), Inches(0.26),
              size=7, color=_SLATE_400)
    if label:
        _add_text(slide, label,
                  _SW - Inches(4.5), _SH - Inches(0.28), Inches(4.3), Inches(0.26),
                  size=7, color=_SLATE_400, align=PP_ALIGN.RIGHT)


def _content_header(slide: Any, title: str, accent: RGBColor = _INDIGO,
                    bg: RGBColor = _SLATE_800) -> None:
    """Full-width dark header band with left accent stripe for content slides."""
    _add_rect(slide, 0, 0, _SW, Inches(0.88), bg)
    _add_rect(slide, 0, 0, Inches(0.35), Inches(0.88), accent)
    _add_text(slide, title,
              Inches(0.55), Inches(0.13), Inches(12.5), Inches(0.62),
              size=18, bold=True, color=_WHITE)


# ── PDF helpers ────────────────────────────────────────────────────────────────

def _make_styles() -> dict[str, ParagraphStyle]:
    base = getSampleStyleSheet()
    return {
        "doc_title": ParagraphStyle(
            "DocTitle", parent=base["Normal"],
            fontSize=22, fontName="Helvetica-Bold",
            textColor=_P_INDIGO_DARK, leading=28,
        ),
        "doc_sub": ParagraphStyle(
            "DocSub", parent=base["Normal"],
            fontSize=12, fontName="Helvetica",
            textColor=_P_SLATE_600, leading=16, spaceAfter=4,
        ),
        "auth_label": ParagraphStyle(
            "AuthLabel", parent=base["Normal"],
            fontSize=10, fontName="Helvetica-Bold",
            textColor=_P_INDIGO, leading=14,
        ),
        "section_title": ParagraphStyle(
            "SectionTitle", parent=base["Normal"],
            fontSize=13, fontName="Helvetica-Bold",
            textColor=_P_SLATE_800, leading=18, spaceBefore=6,
        ),
        "summary_body": ParagraphStyle(
            "SummaryBody", parent=base["Normal"],
            fontSize=10, fontName="Helvetica",
            textColor=_P_SLATE_600, leading=16,
        ),
        "gap_title": ParagraphStyle(
            "GapTitle", parent=base["Normal"],
            fontSize=11, fontName="Helvetica-Bold",
            textColor=_P_SLATE_800, leading=15,
        ),
        "gap_label": ParagraphStyle(
            "GapLabel", parent=base["Normal"],
            fontSize=8, fontName="Helvetica-Bold",
            textColor=_P_INDIGO, leading=11,
        ),
        "gap_body": ParagraphStyle(
            "GapBody", parent=base["Normal"],
            fontSize=9, fontName="Helvetica",
            textColor=_P_SLATE_600, leading=14,
        ),
        "action_title": ParagraphStyle(
            "ActionTitle", parent=base["Normal"],
            fontSize=10, fontName="Helvetica-Bold",
            textColor=_P_SLATE_800, leading=14,
        ),
        "action_body": ParagraphStyle(
            "ActionBody", parent=base["Normal"],
            fontSize=9, fontName="Helvetica",
            textColor=_P_SLATE_600, leading=13,
        ),
        "footer": ParagraphStyle(
            "Footer", parent=base["Normal"],
            fontSize=8, fontName="Helvetica",
            textColor=_P_SLATE_400, alignment=TA_CENTER, leading=12,
        ),
    }


def _header_banner(filename: str, page_w: float, st: dict) -> list:
    """Full-width cover banner for the PDF."""
    title_para = Paragraph("Regulatory Analysis Report", st["doc_title"])
    file_para  = Paragraph(filename, st["doc_sub"])
    data = [[title_para], [file_para]]
    t = Table(data, colWidths=[page_w])
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), _P_INDIGO_LIGHT),
        ("TOPPADDING",    (0, 0), (-1, -1), 18),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("LEFTPADDING",   (0, 0), (-1, -1), 20),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 20),
        ("LINEBELOW",     (0, -1), (-1, -1), 2.5, _P_INDIGO),
    ]))
    return [t, Spacer(1, 16)]


def _auth_header(auth: str, page_w: float, st: dict) -> list:
    data = [[Paragraph(f"Authority Context: {auth}", st["auth_label"])]]
    t = Table(data, colWidths=[page_w])
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), _P_INDIGO_LIGHT),
        ("TOPPADDING",    (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING",   (0, 0), (-1, -1), 14),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 14),
        ("LINEAFTER",     (0, 0), (0, -1),  3, _P_INDIGO),
    ]))
    return [t, Spacer(1, 10)]


def _summary_card(summary: str, page_w: float, st: dict) -> list:
    label = Paragraph("EXECUTIVE SUMMARY", st["gap_label"])
    body  = Paragraph(_sanitize(summary), st["summary_body"])
    data  = [[label], [body]]
    t = Table(data, colWidths=[page_w - 10])
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), _P_SLATE_50),
        ("TOPPADDING",    (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING",   (0, 0), (-1, -1), 14),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 14),
        ("BOX",           (0, 0), (-1, -1), 0.5, _P_BORDER),
        ("LINEAFTER",     (0, 0), (0, -1),  3, _P_INDIGO),
    ]))
    return [t, Spacer(1, 14)]


def _gap_card(gap: dict, page_w: float, st: dict) -> Any:
    sev      = gap.get("severity", "Low")
    sev_c    = _P_SEV.get(sev, _P_SEV["Low"])
    sev_bg   = _P_SEV_BG.get(sev, _P_SEV_BG["Low"])
    bar_w    = 5
    cont_w   = page_w - bar_w - 10
    sev_tag  = Paragraph(
        f'<font color="{sev_c.hexval() if hasattr(sev_c,"hexval") else "#059669"}"><b>{sev.upper()}</b></font>',
        st["gap_label"],
    )
    title_p  = Paragraph(_sanitize(gap.get("title", "")), st["gap_title"])
    desc_p   = Paragraph(
        f'<i>Observation:</i> {_sanitize(gap.get("description", ""))}', st["gap_body"]
    )
    action_p = Paragraph(
        f'<font color="#4338CA">&#x2192; Strategic Action:</font> {_sanitize(gap.get("recommended_action", ""))}',
        st["gap_body"],
    )
    # Title row: title left, severity tag right
    title_row = Table(
        [[title_p, sev_tag]],
        colWidths=[cont_w - 70, 70],
    )
    title_row.setStyle(TableStyle([
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING",    (0, 0), (-1, -1), 0),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
        ("LEFTPADDING",   (0, 0), (-1, -1), 0),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 0),
        ("BACKGROUND",    (1, 0), (1, 0), sev_bg),
        ("BOX",           (1, 0), (1, 0), 0.5, sev_c),
        ("ALIGN",         (1, 0), (1, 0), "CENTER"),
    ]))

    content = [title_row, Spacer(1, 5), desc_p, Spacer(1, 4), action_p]
    data = [["", content]]
    t = Table(data, colWidths=[bar_w, cont_w])
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (0, -1), sev_c),
        ("BACKGROUND",    (1, 0), (1, -1), _P_WHITE),
        ("BOX",           (1, 0), (1, -1), 0.5, _P_BORDER),
        ("TOPPADDING",    (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING",   (1, 0), (1, -1),  12),
        ("RIGHTPADDING",  (1, 0), (1, -1),  12),
        ("LEFTPADDING",   (0, 0), (0, -1),  0),
        ("RIGHTPADDING",  (0, 0), (0, -1),  0),
        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
    ]))
    return KeepTogether([t, Spacer(1, 7)])


def _action_card(action: dict, idx: int, page_w: float, st: dict) -> Any:
    pri_map  = {"High": _P_SEV["Critical"], "Medium": _P_SEV["Medium"], "Low": _P_SEV["Low"]}
    pri      = action.get("priority", "Medium")
    pri_c    = pri_map.get(pri, _P_SEV["Low"])
    num_para = Paragraph(f"<b>{idx + 1}</b>", ParagraphStyle(
        "Num", parent=st["action_title"],
        textColor=pri_c, alignment=TA_CENTER, fontSize=13,
    ))
    title_p  = Paragraph(_sanitize(action.get("title", "")), st["action_title"])
    desc_p   = Paragraph(_sanitize(action.get("description", "")), st["action_body"])
    pri_para = Paragraph(
        f'<font color="{pri_c.hexval() if hasattr(pri_c,"hexval") else "#EF4444"}"><b>{pri} Priority</b></font>',
        st["gap_label"],
    )
    data = [[num_para, [pri_para, Spacer(1, 3), title_p, Spacer(1, 3), desc_p]]]
    t = Table(data, colWidths=[30, page_w - 40])
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), _P_SLATE_50),
        ("BOX",           (0, 0), (-1, -1), 0.5, _P_BORDER),
        ("TOPPADDING",    (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING",   (0, 0), (-1, -1), 10),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 10),
        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
        ("LINEAFTER",     (0, 0), (0, -1), 1, _P_BORDER),
    ]))
    return KeepTogether([t, Spacer(1, 6)])


def _section_header(label: str, page_w: float, st: dict) -> list:
    data = [[Paragraph(label, st["section_title"])]]
    t = Table(data, colWidths=[page_w])
    t.setStyle(TableStyle([
        ("LINEBELOW",     (0, 0), (-1, -1), 1.5, _P_INDIGO),
        ("TOPPADDING",    (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("LEFTPADDING",   (0, 0), (-1, -1), 0),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 0),
    ]))
    return [t, Spacer(1, 8)]


class ExportService:
    # ── PDF ───────────────────────────────────────────────────────────────────
    async def generate_pdf(self, analysis_data: dict[str, Any], filename: str) -> str:
        buffer = io.BytesIO()
        margin = 0.75 * inch
        doc    = SimpleDocTemplate(
            buffer, pagesize=letter,
            leftMargin=margin, rightMargin=margin,
            topMargin=0.75 * inch, bottomMargin=0.75 * inch,
        )
        page_w = letter[0] - 2 * margin
        st     = _make_styles()
        elems: list = []

        elems += _header_banner(_sanitize(filename), page_w, st)

        for auth, data in analysis_data.items():
            elems += _auth_header(auth, page_w, st)
            elems += _summary_card(data.get("summary", "N/A"), page_w, st)

            gaps    = data.get("gaps", [])
            actions = data.get("actions", [])

            if gaps:
                elems += _section_header("Detailed Findings & Regulatory Risks", page_w, st)
                for gap in gaps:
                    elems.append(_gap_card(gap, page_w, st))
                elems.append(Spacer(1, 8))

            if actions:
                elems += _section_header("Strategic Recommendations & Next Steps", page_w, st)
                for i, act in enumerate(actions):
                    elems.append(_action_card(act, i, page_w, st))
            else:
                elems += _section_header("Strategic Recommendations & Next Steps", page_w, st)
                elems.append(Paragraph(
                    "Based on the analysis, it is recommended to prioritize 'Critical' gaps "
                    "before the next submission milestone. Ensure all manufacturing process "
                    "validation data is synchronized with the clinical study protocol.",
                    st["summary_body"],
                ))

            elems.append(Spacer(1, 20))
            elems.append(HRFlowable(width=page_w, thickness=0.5, color=_P_BORDER))
            elems.append(Spacer(1, 20))

        elems.append(Paragraph(
            "This report is AI-generated and for preparation purposes only. "
            "Gulfstream Intelligence · Confidential",
            st["footer"],
        ))

        doc.build(elems)
        url = await storage_service.upload_file(
            buffer.getvalue(), f"Report_{filename}.pdf", "application/pdf"
        )
        return url

    # ── DOCX ──────────────────────────────────────────────────────────────────
    async def generate_docx(self, analysis_data: dict[str, Any], filename: str) -> str:

        # ── Colour constants ─────────────────────────────────────────────────
        _C = {
            "indigo":       "6366F1",
            "indigo_dark":  "4338CA",
            "indigo_light": "EEF2FF",
            "slate_800":    "1E293B",
            "slate_600":    "475569",
            "slate_400":    "94A3B8",
            "white":        "FFFFFF",
            "bg":           "F8FAFC",
            "border":       "E2E8F0",
            "purple":       "7C3AED",
            "purple_light": "FAF5FF",
            "green":        "059669",
            "green_bg":     "F0FDF4",
            "green_dark":   "10B981",
            "red":          "EF4444",
            "red_bg":       "FEF2F2",
        }
        _SEV_C  = {"Critical": "EF4444", "High": "F97316", "Medium": "D97706", "Low": "059669"}
        _SEV_BG = {"Critical": "FEF2F2", "High": "FFF7ED", "Medium": "FFFBEB", "Low": "F0FDF4"}
        _PRI_C  = {"High": "EF4444", "Medium": "D97706", "Low": "059669"}
        _PRI_BG = {"High": "FEF2F2", "Medium": "FFFBEB", "Low": "F0FDF4"}

        def _rgb(h: str) -> DocxRGB:
            h = h.lstrip("#")
            return DocxRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def _cell_bg(cell: Any, hex_color: str) -> None:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), hex_color.lstrip("#"))
            tcPr.append(shd)

        def _cell_borders(cell: Any, sides: dict) -> None:
            """sides: {side_name: hex_color}. side_name ∈ top|bottom|left|right"""
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement("w:tcBorders")
            for side, color in sides.items():
                el = OxmlElement(f"w:{side}")
                el.set(qn("w:val"), "single")
                el.set(qn("w:sz"), "8")
                el.set(qn("w:color"), color.lstrip("#"))
                tcBorders.append(el)
            tcPr.append(tcBorders)

        def _no_borders(table: Any) -> None:
            t = table._tbl
            tblPr = t.tblPr
            tblBorders = OxmlElement("w:tblBorders")
            for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
                el = OxmlElement(f"w:{side}")
                el.set(qn("w:val"), "none")
                el.set(qn("w:sz"), "0")
                el.set(qn("w:color"), "auto")
                tblBorders.append(el)
            existing = tblPr.find(qn("w:tblBorders"))
            if existing is not None:
                tblPr.remove(existing)
            tblPr.append(tblBorders)

        def _set_col_widths(table: Any, *widths_in: float) -> None:
            table.autofit = False
            for col, w in zip(table.columns, widths_in):
                for cell in col.cells:
                    cell.width = DocxInches(w)

        def _run(para: Any, text: str, size: int = 10, bold: bool = False,
                 italic: bool = False, color: str = "1E293B") -> None:
            r = para.add_run(text)
            r.font.size = DocxPt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = _rgb(color)

        def _para(cell_or_doc: Any, size: int = 10, bold: bool = False,
                  italic: bool = False, color: str = "1E293B",
                  sb: int = 0, sa: int = 6, indent: float = 0,
                  align: int = WD_ALIGN_PARAGRAPH.LEFT) -> Any:
            """Return a fresh paragraph appended to cell_or_doc."""
            if hasattr(cell_or_doc, "add_paragraph"):
                p = cell_or_doc.add_paragraph()
            else:
                # first paragraph of a fresh cell
                p = cell_or_doc.paragraphs[0]
            p.alignment = align
            p.paragraph_format.space_before = DocxPt(sb)
            p.paragraph_format.space_after  = DocxPt(sa)
            if indent:
                p.paragraph_format.left_indent = DocxInches(indent)
            return p

        def _section_heading(doc: Any, label: str, stripe: str, stripe_bg: str) -> None:
            """Coloured section heading bar."""
            tbl = doc.add_table(rows=1, cols=1)
            _no_borders(tbl)
            c = tbl.cell(0, 0)
            _cell_bg(c, stripe_bg)
            _cell_borders(c, {"left": stripe, "bottom": stripe})
            p = c.paragraphs[0]
            p.paragraph_format.space_before = DocxPt(8)
            p.paragraph_format.space_after  = DocxPt(8)
            p.paragraph_format.left_indent  = DocxInches(0.12)
            _run(p, label, size=9, bold=True, color=stripe)
            doc.add_paragraph().paragraph_format.space_after = DocxPt(2)

        def _spacer(doc: Any, pts: int = 6) -> None:
            sp = doc.add_paragraph()
            sp.paragraph_format.space_before = DocxPt(0)
            sp.paragraph_format.space_after  = DocxPt(pts)

        def _hr(doc: Any) -> None:
            p = doc.add_paragraph()
            p.paragraph_format.space_before = DocxPt(10)
            p.paragraph_format.space_after  = DocxPt(10)
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bt = OxmlElement("w:bottom")
            bt.set(qn("w:val"), "single")
            bt.set(qn("w:sz"), "6")
            bt.set(qn("w:color"), _C["border"])
            pBdr.append(bt)
            pPr.append(pBdr)

        # ── Document setup ────────────────────────────────────────────────────
        doc = DocxDocument()
        for section in doc.sections:
            section.top_margin    = DocxInches(0.75)
            section.bottom_margin = DocxInches(0.75)
            section.left_margin   = DocxInches(0.75)
            section.right_margin  = DocxInches(0.75)

        # ── Cover banner ──────────────────────────────────────────────────────
        banner = doc.add_table(rows=3, cols=1)
        _no_borders(banner)

        # Thin top accent stripe
        top_cell = banner.cell(0, 0)
        _cell_bg(top_cell, _C["indigo"])
        tp = top_cell.paragraphs[0]
        tp.paragraph_format.space_before = DocxPt(3)
        tp.paragraph_format.space_after  = DocxPt(3)
        _run(tp, " ", size=3, color=_C["indigo"])

        # Title row
        title_cell = banner.cell(1, 0)
        _cell_bg(title_cell, _C["indigo_dark"])
        title_p = title_cell.paragraphs[0]
        title_p.paragraph_format.space_before = DocxPt(16)
        title_p.paragraph_format.space_after  = DocxPt(8)
        title_p.paragraph_format.left_indent  = DocxInches(0.25)
        _run(title_p, "REGULATORY ANALYSIS REPORT", size=22, bold=True, color=_C["white"])

        # Subtitle row
        sub_cell = banner.cell(2, 0)
        _cell_bg(sub_cell, _C["indigo_light"])
        sub_p = sub_cell.paragraphs[0]
        sub_p.paragraph_format.space_before = DocxPt(10)
        sub_p.paragraph_format.space_after  = DocxPt(4)
        sub_p.paragraph_format.left_indent  = DocxInches(0.25)
        _run(sub_p, f"Document:  ", size=11, bold=True, color=_C["indigo_dark"])
        _run(sub_p, filename, size=11, color=_C["indigo_dark"])
        sub_p2 = sub_cell.add_paragraph()
        sub_p2.paragraph_format.space_before = DocxPt(0)
        sub_p2.paragraph_format.space_after  = DocxPt(10)
        sub_p2.paragraph_format.left_indent  = DocxInches(0.25)
        _run(sub_p2, "Prepared by Gulfstream Intelligence  ·  Confidential",
             size=9, italic=True, color=_C["slate_600"])

        _spacer(doc, 10)

        # ── Per-authority sections ─────────────────────────────────────────────
        for auth, data in analysis_data.items():
            summary    = data.get("summary", "N/A")
            gaps       = data.get("gaps", [])
            insights   = data.get("insights", [])
            risks      = data.get("risks", [])
            actions    = data.get("actions", [])
            confidence = data.get("confidence_score", 0)

            # Authority header bar (two-column: label | confidence)
            auth_tbl = doc.add_table(rows=1, cols=2)
            _no_borders(auth_tbl)
            _set_col_widths(auth_tbl, 5.8, 1.2)

            auth_c = auth_tbl.cell(0, 0)
            _cell_bg(auth_c, _C["indigo"])
            ap = auth_c.paragraphs[0]
            ap.paragraph_format.space_before = DocxPt(10)
            ap.paragraph_format.space_after  = DocxPt(10)
            ap.paragraph_format.left_indent  = DocxInches(0.18)
            _run(ap, f"Authority: {auth}", size=13, bold=True, color=_C["white"])

            conf_c = auth_tbl.cell(0, 1)
            _cell_bg(conf_c, _C["indigo_dark"])
            cp = conf_c.paragraphs[0]
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cp.paragraph_format.space_before = DocxPt(10)
            cp.paragraph_format.space_after  = DocxPt(10)
            _run(cp, f"{round(confidence * 100)}%", size=13, bold=True, color=_C["white"])
            cp2 = conf_c.add_paragraph()
            cp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cp2.paragraph_format.space_before = DocxPt(0)
            cp2.paragraph_format.space_after  = DocxPt(10)
            _run(cp2, "confidence", size=8, color="C7D2FE")

            _spacer(doc, 6)

            # Executive Summary card
            sum_tbl = doc.add_table(rows=2, cols=1)
            _no_borders(sum_tbl)

            label_c = sum_tbl.cell(0, 0)
            _cell_bg(label_c, _C["indigo_light"])
            _cell_borders(label_c, {"left": _C["indigo"]})
            lp = label_c.paragraphs[0]
            lp.paragraph_format.space_before = DocxPt(7)
            lp.paragraph_format.space_after  = DocxPt(7)
            lp.paragraph_format.left_indent  = DocxInches(0.14)
            _run(lp, "EXECUTIVE SUMMARY", size=8, bold=True, color=_C["indigo_dark"])

            body_c = sum_tbl.cell(1, 0)
            _cell_bg(body_c, _C["bg"])
            _cell_borders(body_c, {"left": _C["indigo"], "bottom": _C["border"]})
            bp = body_c.paragraphs[0]
            bp.paragraph_format.space_before = DocxPt(8)
            bp.paragraph_format.space_after  = DocxPt(10)
            bp.paragraph_format.left_indent  = DocxInches(0.14)
            _run(bp, summary, size=10, color=_C["slate_600"])

            _spacer(doc, 8)

            # ── Key Insights ──────────────────────────────────────────────────
            if insights:
                _section_heading(doc, "KEY INSIGHTS", _C["purple"], _C["purple_light"])
                for i, ins in enumerate(insights):
                    ins_tbl = doc.add_table(rows=1, cols=2)
                    _no_borders(ins_tbl)
                    _set_col_widths(ins_tbl, 0.38, 6.62)

                    num_c = ins_tbl.cell(0, 0)
                    _cell_bg(num_c, _C["purple"])
                    np = num_c.paragraphs[0]
                    np.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    np.paragraph_format.space_before = DocxPt(8)
                    np.paragraph_format.space_after  = DocxPt(8)
                    _run(np, str(i + 1), size=11, bold=True, color=_C["white"])

                    txt_c = ins_tbl.cell(0, 1)
                    _cell_bg(txt_c, _C["purple_light"])
                    _cell_borders(txt_c, {"bottom": "E9D5FF"})
                    tp2 = txt_c.paragraphs[0]
                    tp2.paragraph_format.space_before = DocxPt(8)
                    tp2.paragraph_format.space_after  = DocxPt(8)
                    tp2.paragraph_format.left_indent  = DocxInches(0.12)
                    _run(tp2, ins.get("content", ""), size=10, color=_C["slate_600"])

                    _spacer(doc, 3)
                _spacer(doc, 4)

            # ── Regulatory Gaps ────────────────────────────────────────────────
            if gaps:
                _section_heading(doc, "REGULATORY GAPS & FINDINGS", _C["red"], _C["red_bg"])
                for gap in gaps:
                    sev  = gap.get("severity", "Low")
                    sc   = _SEV_C.get(sev, _SEV_C["Low"])
                    sbg  = _SEV_BG.get(sev, _SEV_BG["Low"])

                    g_tbl = doc.add_table(rows=1, cols=2)
                    _no_borders(g_tbl)
                    _set_col_widths(g_tbl, 0.12, 6.88)

                    bar_c = g_tbl.cell(0, 0)
                    _cell_bg(bar_c, sc)
                    bar_c.paragraphs[0].paragraph_format.space_before = DocxPt(8)

                    con_c = g_tbl.cell(0, 1)
                    _cell_bg(con_c, _C["white"])
                    _cell_borders(con_c, {"top": _C["border"], "right": _C["border"], "bottom": _C["border"]})

                    g_title = con_c.paragraphs[0]
                    g_title.paragraph_format.space_before = DocxPt(8)
                    g_title.paragraph_format.space_after  = DocxPt(3)
                    g_title.paragraph_format.left_indent  = DocxInches(0.12)
                    _run(g_title, gap.get("title", ""), size=11, bold=True, color=_C["slate_800"])
                    _run(g_title, f"  ·  {sev}", size=9, bold=True, color=sc)

                    g_desc = con_c.add_paragraph()
                    g_desc.paragraph_format.space_before = DocxPt(0)
                    g_desc.paragraph_format.space_after  = DocxPt(3)
                    g_desc.paragraph_format.left_indent  = DocxInches(0.12)
                    _run(g_desc, "Observation: ", size=9, italic=True, color=_C["slate_400"])
                    _run(g_desc, gap.get("description", ""), size=9, color=_C["slate_600"])

                    if gap.get("regulatory_impact"):
                        g_impact = con_c.add_paragraph()
                        g_impact.paragraph_format.space_before = DocxPt(0)
                        g_impact.paragraph_format.space_after  = DocxPt(3)
                        g_impact.paragraph_format.left_indent  = DocxInches(0.12)
                        _run(g_impact, "Impact: ", size=9, italic=True, color=_C["slate_400"])
                        _run(g_impact, gap.get("regulatory_impact", ""), size=9, color=_C["slate_600"])

                    g_act = con_c.add_paragraph()
                    g_act.paragraph_format.space_before = DocxPt(2)
                    g_act.paragraph_format.space_after  = DocxPt(9)
                    g_act.paragraph_format.left_indent  = DocxInches(0.12)
                    _run(g_act, "→ Strategic Action: ", size=9, bold=True, color=_C["indigo_dark"])
                    _run(g_act, gap.get("recommended_action", ""), size=9, color=_C["indigo_dark"])

                    _spacer(doc, 3)
                _spacer(doc, 4)

            # ── Key Risks ──────────────────────────────────────────────────────
            if risks:
                _section_heading(doc, "KEY RISKS", _C["red"], _C["red_bg"])
                for risk in risks:
                    r_tbl = doc.add_table(rows=1, cols=2)
                    _no_borders(r_tbl)
                    _set_col_widths(r_tbl, 0.28, 6.72)

                    bullet_c = r_tbl.cell(0, 0)
                    _cell_bg(bullet_c, _C["red_bg"])
                    bpp = bullet_c.paragraphs[0]
                    bpp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    bpp.paragraph_format.space_before = DocxPt(7)
                    bpp.paragraph_format.space_after  = DocxPt(7)
                    _run(bpp, "!", size=11, bold=True, color=_C["red"])

                    risk_c = r_tbl.cell(0, 1)
                    _cell_bg(risk_c, _C["red_bg"])
                    _cell_borders(risk_c, {"bottom": "FECACA"})
                    rp2 = risk_c.paragraphs[0]
                    rp2.paragraph_format.space_before = DocxPt(7)
                    rp2.paragraph_format.space_after  = DocxPt(7)
                    rp2.paragraph_format.left_indent  = DocxInches(0.08)
                    _run(rp2, risk, size=10, color=_C["slate_600"])

                    _spacer(doc, 3)
                _spacer(doc, 4)

            # ── Strategic Recommendations ──────────────────────────────────────
            if actions:
                _section_heading(doc, "STRATEGIC RECOMMENDATIONS", _C["green"], _C["green_bg"])
                for i, action in enumerate(actions):
                    pri  = action.get("priority", "Medium")
                    pc   = _PRI_C.get(pri, _PRI_C["Medium"])
                    pbg  = _PRI_BG.get(pri, _PRI_BG["Medium"])

                    a_tbl = doc.add_table(rows=1, cols=2)
                    _no_borders(a_tbl)
                    _set_col_widths(a_tbl, 0.42, 6.58)

                    num_c2 = a_tbl.cell(0, 0)
                    _cell_bg(num_c2, _C["green_dark"])
                    np2 = num_c2.paragraphs[0]
                    np2.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    np2.paragraph_format.space_before = DocxPt(10)
                    np2.paragraph_format.space_after  = DocxPt(10)
                    _run(np2, str(i + 1), size=13, bold=True, color=_C["white"])

                    con_c2 = a_tbl.cell(0, 1)
                    _cell_bg(con_c2, _C["bg"])
                    _cell_borders(con_c2, {"top": _C["border"], "right": _C["border"], "bottom": _C["border"]})

                    a_title = con_c2.paragraphs[0]
                    a_title.paragraph_format.space_before = DocxPt(8)
                    a_title.paragraph_format.space_after  = DocxPt(2)
                    a_title.paragraph_format.left_indent  = DocxInches(0.12)
                    _run(a_title, action.get("title", ""), size=11, bold=True, color=_C["slate_800"])
                    _run(a_title, f"  ·  {pri} Priority", size=9, bold=True, color=pc)

                    a_desc = con_c2.add_paragraph()
                    a_desc.paragraph_format.space_before = DocxPt(0)
                    a_desc.paragraph_format.space_after  = DocxPt(9)
                    a_desc.paragraph_format.left_indent  = DocxInches(0.12)
                    _run(a_desc, action.get("description", ""), size=10, color=_C["slate_600"])

                    _spacer(doc, 3)

            _hr(doc)

        # ── Footer ────────────────────────────────────────────────────────────
        _spacer(doc, 8)
        footer = doc.add_paragraph()
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer.paragraph_format.space_before = DocxPt(6)
        _run(footer,
             "This report is AI-generated and for preparation purposes only.  "
             "Gulfstream Intelligence  ·  Confidential",
             size=8, italic=True, color=_C["slate_400"])

        buf = io.BytesIO()
        doc.save(buf)
        url = await storage_service.upload_file(
            buf.getvalue(),
            f"Report_{filename}.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        return url

    # ── PPTX ─────────────────────────────────────────────────────────────────
    async def generate_pptx(self, analysis_data: dict[str, Any], filename: str) -> str:
        prs = Presentation()
        prs.slide_width  = _SW
        prs.slide_height = _SH

        # ── Cover slide (dramatic dark) ───────────────────────────────────────
        slide = _blank(prs)
        _solid_bg(slide, _SLATE_800)

        # Decorative stacked geometry — top-right corner
        _add_rect(slide, _SW - Inches(5.2), 0, Inches(5.2), Inches(3.2), _INDIGO_DARK)
        _add_rect(slide, _SW - Inches(3.8), 0, Inches(3.8), Inches(2.2), _INDIGO)
        _add_rect(slide, _SW - Inches(2.2), 0, Inches(2.2), Inches(1.2), _INDIGO_MID)

        # Left accent stripe
        _add_rect(slide, 0, 0, Inches(0.38), _SH, _INDIGO)

        # Bottom band
        _add_rect(slide, Inches(0.38), _SH - Inches(1.3), _SW - Inches(0.38), Inches(1.3), _INDIGO_DARK)

        # Title block
        _add_text(slide, "REGULATORY",
                  Inches(0.72), Inches(1.55), Inches(8.5), Inches(0.55),
                  size=16, bold=True, color=_SLATE_400)
        _add_text(slide, "ANALYSIS REPORT",
                  Inches(0.72), Inches(2.1), Inches(8.5), Inches(1.1),
                  size=38, bold=True, color=_WHITE)
        # Accent divider line
        _add_rect(slide, Inches(0.72), Inches(3.3), Inches(4.2), Inches(0.05), _INDIGO)

        _add_text(slide, filename,
                  Inches(0.72), Inches(3.48), Inches(8.5), Inches(0.55),
                  size=13, color=_SLATE_400)
        _add_text(slide, "Gulfstream Intelligence  ·  Regulatory Intelligence Platform",
                  Inches(0.72), Inches(4.1), Inches(8.5), Inches(0.42),
                  size=10, color=_SLATE_400)

        _add_text(slide, "Confidential  ·  For preparation purposes only.",
                  Inches(0.72), _SH - Inches(1.18), Inches(9.5), Inches(0.38),
                  size=9, color=_SLATE_600)

        for auth, data in analysis_data.items():
            summary    = data.get("summary", "")
            gaps       = data.get("gaps", [])
            insights   = data.get("insights", [])
            actions    = data.get("actions", [])
            risks      = data.get("risks", [])
            confidence = data.get("confidence_score", 0)

            # ── Authority divider slide ───────────────────────────────────────
            slide = _blank(prs)
            _solid_bg(slide, _INDIGO)
            # Right dark panel
            _add_rect(slide, _SW - Inches(4.8), 0, Inches(4.8), _SH, _INDIGO_DARK)
            # Decorative accent bar across middle
            _add_rect(slide, 0, _SH / 2 - Inches(0.04), _SW - Inches(4.8), Inches(0.04),
                      _INDIGO_MID)

            _add_text(slide, "AUTHORITY",
                      Inches(0.55), Inches(1.7), Inches(8.0), Inches(0.5),
                      size=12, bold=True, color=_INDIGO_LIGHT)
            _add_text(slide, auth.upper(),
                      Inches(0.55), Inches(2.15), Inches(8.0), Inches(1.4),
                      size=46, bold=True, color=_WHITE)
            _add_text(slide, "Section Overview  ·  Regulatory Analysis",
                      Inches(0.55), Inches(3.65), Inches(8.0), Inches(0.5),
                      size=13, color=_INDIGO_LIGHT)

            # Confidence badge in right panel
            _add_text(slide, "CONFIDENCE",
                      _SW - Inches(4.3), Inches(2.5), Inches(4.0), Inches(0.4),
                      size=10, bold=True, color=_SLATE_400, align=PP_ALIGN.CENTER)
            _add_text(slide, f"{round(confidence * 100)}%",
                      _SW - Inches(4.3), Inches(2.9), Inches(4.0), Inches(1.4),
                      size=52, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

            # ── Executive Summary slide (two-column) ──────────────────────────
            slide = _blank(prs)
            _solid_bg(slide, _WHITE)
            _content_header(slide, f"{auth}  ·  Executive Summary")

            # Left text panel
            _add_rect(slide, Inches(0.4), Inches(0.98), Inches(7.8), Inches(6.1),
                      RGBColor(0xF8, 0xFA, 0xFF), line=RGBColor(0xE2, 0xE8, 0xF0))
            _add_rect(slide, Inches(0.4), Inches(0.98), Inches(0.2), Inches(6.1), _INDIGO)
            _add_text(slide, "EXECUTIVE SUMMARY",
                      Inches(0.78), Inches(1.08), Inches(7.1), Inches(0.32),
                      size=7, bold=True, color=_SLATE_400)
            _add_text(slide, summary,
                      Inches(0.78), Inches(1.44), Inches(7.1), Inches(5.5),
                      size=11, color=_SLATE_600)

            # Right stats panel
            rx = Inches(8.5)
            rw = _SW - rx - Inches(0.35)

            # Confidence badge
            _add_rect(slide, rx, Inches(0.98), rw, Inches(2.5), _INDIGO_DARK)
            _add_text(slide, "CONFIDENCE SCORE",
                      rx + Inches(0.2), Inches(1.12), rw - Inches(0.3), Inches(0.32),
                      size=8, bold=True, color=_INDIGO_LIGHT)
            _add_text(slide, f"{round(confidence * 100)}%",
                      rx, Inches(1.45), rw, Inches(1.45),
                      size=46, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

            # Stats: gaps / insights / actions
            stats = [("Gaps", str(len(gaps))), ("Insights", str(len(insights))), ("Actions", str(len(actions)))]
            stat_w = rw / 3
            for si, (lbl, val) in enumerate(stats):
                sx = rx + si * stat_w
                _add_rect(slide, sx, Inches(3.58), stat_w - Inches(0.06), Inches(1.95),
                          _INDIGO_LIGHT, line=RGBColor(0xC7, 0xD2, 0xFE))
                _add_text(slide, val,
                          sx, Inches(3.72), stat_w - Inches(0.06), Inches(1.0),
                          size=30, bold=True, color=_INDIGO_DARK, align=PP_ALIGN.CENTER)
                _add_text(slide, lbl,
                          sx, Inches(4.75), stat_w - Inches(0.06), Inches(0.45),
                          size=9, color=_SLATE_600, align=PP_ALIGN.CENTER)

            _slide_footer(slide, f"Authority: {auth}")

            # ── Key Insights slide (purple theme) ─────────────────────────────
            if insights:
                INSIGHTS_PER = 6
                for si in range(0, len(insights), INSIGHTS_PER):
                    batch_ins = insights[si: si + INSIGHTS_PER]
                    slide = _blank(prs)
                    _solid_bg(slide, _WHITE)
                    _content_header(slide, f"{auth}  ·  Key Insights  ({len(insights)} total)",
                                    accent=_PURPLE, bg=_PURPLE_DARK)

                    cols, rows = 2, 3
                    card_w = Inches(6.2)
                    card_h = Inches(1.88)
                    gap_x, gap_y = Inches(0.4), Inches(0.12)

                    for idx, ins in enumerate(batch_ins):
                        col = idx % cols
                        row = idx // cols
                        x = Inches(0.35) + col * (card_w + gap_x)
                        y = Inches(1.02) + row * (card_h + gap_y)

                        _add_rect(slide, x, y, card_w, card_h,
                                  _PURPLE_LIGHT, line=RGBColor(0xE9, 0xD5, 0xFF))
                        _add_rect(slide, x, y, Inches(0.42), card_h, _PURPLE)
                        # Number
                        _add_text(slide, str(si + idx + 1),
                                  x + Inches(0.02), y + Inches(0.55), Inches(0.38), Inches(0.75),
                                  size=16, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
                        _add_text(slide, ins.get("content", "")[:185],
                                  x + Inches(0.52), y + Inches(0.12), card_w - Inches(0.65), card_h - Inches(0.22),
                                  size=10, color=_SLATE_600)

                    _slide_footer(slide, f"Authority: {auth}")

            # ── Gaps slides ───────────────────────────────────────────────────
            GAPS_PER = 3
            for si in range(0, len(gaps), GAPS_PER):
                batch = gaps[si: si + GAPS_PER]
                slide = _blank(prs)
                _solid_bg(slide, _WHITE)
                pg = si // GAPS_PER + 1
                pages = (len(gaps) + GAPS_PER - 1) // GAPS_PER
                _content_header(
                    slide,
                    f"{auth}  ·  Regulatory Gaps  —  {len(gaps)} found  ({pg}/{pages})",
                    accent=_RED,
                )

                row_h   = Inches(2.05)
                row_top = Inches(1.0)

                for i, gap in enumerate(batch):
                    top = row_top + i * row_h
                    sev = gap.get("severity", "Low")
                    sc  = _SEV_COLOR.get(sev, _GREEN)
                    sbg = _SEV_BG_COLOR.get(sev, _GREEN_LIGHT)

                    # Card
                    _add_rect(slide, Inches(0.4), top, _SW - Inches(0.8), row_h - Inches(0.1),
                              RGBColor(0xF8, 0xFA, 0xFF), line=RGBColor(0xE2, 0xE8, 0xF0))
                    # Severity stripe
                    _add_rect(slide, Inches(0.4), top, Inches(0.22), row_h - Inches(0.1), sc)
                    # Severity pill
                    _add_rect(slide, Inches(0.78), top + Inches(0.18), Inches(1.05), Inches(0.32),
                              sbg, line=sc)
                    _add_text(slide, sev.upper(),
                              Inches(0.80), top + Inches(0.19), Inches(1.01), Inches(0.28),
                              size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)
                    # Title
                    _add_text(slide, gap.get("title", ""),
                              Inches(2.0), top + Inches(0.12), Inches(11.0), Inches(0.44),
                              size=13, bold=True, color=_SLATE_800)
                    # Description
                    desc = gap.get("description", "")
                    _add_text(slide, desc[:190] + ("…" if len(desc) > 190 else ""),
                              Inches(0.8), top + Inches(0.62), Inches(12.15), Inches(0.54),
                              size=10, color=_SLATE_600)
                    # Action
                    act_txt = gap.get("recommended_action", "")
                    _add_text(slide, "→  " + (act_txt[:165] + ("…" if len(act_txt) > 165 else "")),
                              Inches(0.8), top + Inches(1.2), Inches(12.15), Inches(0.52),
                              size=10, color=_INDIGO_DARK)

                _slide_footer(slide, f"Authority: {auth}")

            # ── Risks slide ───────────────────────────────────────────────────
            if risks:
                slide = _blank(prs)
                _solid_bg(slide, _WHITE)
                _content_header(slide, f"{auth}  ·  Key Risks  —  {len(risks)} identified",
                                accent=_RED, bg=_RED_DARK)

                for idx, risk in enumerate(risks[:6]):
                    top = Inches(1.02) + idx * Inches(1.06)
                    _add_rect(slide, Inches(0.4), top, _SW - Inches(0.8), Inches(0.92),
                              _RED_LIGHT, line=RGBColor(0xFE, 0xCA, 0xCA))
                    _add_rect(slide, Inches(0.4), top, Inches(0.22), Inches(0.92), _RED)
                    _add_text(slide, f"R{idx + 1}",
                              Inches(0.44), top + Inches(0.18), Inches(0.2), Inches(0.56),
                              size=10, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
                    _add_text(slide, risk[:230] + ("…" if len(risk) > 230 else ""),
                              Inches(0.82), top + Inches(0.1), Inches(12.15), Inches(0.72),
                              size=11, color=_SLATE_600)

                _slide_footer(slide, f"Authority: {auth}")

            # ── Recommendations slides ────────────────────────────────────────
            if actions:
                ACTS_PER = 4
                pri_c_map  = {"High": _RED,    "Medium": _YELLOW,    "Low": _GREEN}
                pri_bg_map = {"High": _RED_LIGHT, "Medium": _YELLOW_LIGHT, "Low": _GREEN_LIGHT}

                for si in range(0, len(actions), ACTS_PER):
                    batch_a = actions[si: si + ACTS_PER]
                    slide = _blank(prs)
                    _solid_bg(slide, _WHITE)
                    pg2 = si // ACTS_PER + 1
                    pages2 = (len(actions) + ACTS_PER - 1) // ACTS_PER
                    label = f"{auth}  ·  Strategic Recommendations  ({pg2}/{pages2})"
                    _content_header(slide, label, accent=_GREEN, bg=_GREEN_DARK2)

                    row_h2  = Inches(1.56)
                    row_top = Inches(1.0)

                    for idx, act in enumerate(batch_a):
                        top = row_top + idx * row_h2
                        p   = act.get("priority", "Medium")
                        pc  = pri_c_map.get(p, _GREEN)
                        pbg = pri_bg_map.get(p, _GREEN_LIGHT)

                        # Card bg
                        _add_rect(slide, Inches(0.4), top, _SW - Inches(0.8), row_h2 - Inches(0.1),
                                  RGBColor(0xF8, 0xFD, 0xF8), line=RGBColor(0xD1, 0xFA, 0xE5))
                        # Numbered left column (green)
                        _add_rect(slide, Inches(0.4), top, Inches(0.55), row_h2 - Inches(0.1), _GREEN)
                        _add_text(slide, str(si + idx + 1),
                                  Inches(0.4), top + Inches(0.4), Inches(0.55), Inches(0.72),
                                  size=16, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
                        # Priority pill
                        _add_rect(slide, Inches(1.1), top + Inches(0.16), Inches(1.1), Inches(0.3),
                                  pbg, line=pc)
                        _add_text(slide, p.upper(),
                                  Inches(1.12), top + Inches(0.17), Inches(1.06), Inches(0.26),
                                  size=8, bold=True, color=pc, align=PP_ALIGN.CENTER)
                        # Title
                        _add_text(slide, act.get("title", ""),
                                  Inches(2.35), top + Inches(0.1), Inches(10.7), Inches(0.44),
                                  size=13, bold=True, color=_SLATE_800)
                        # Description
                        desc = act.get("description", "")
                        _add_text(slide, desc[:195] + ("…" if len(desc) > 195 else ""),
                                  Inches(1.1), top + Inches(0.55), Inches(12.0), Inches(0.85),
                                  size=10, color=_SLATE_600)

                    _slide_footer(slide, f"Authority: {auth}")

        # ── Closing slide ─────────────────────────────────────────────────────
        slide = _blank(prs)
        _solid_bg(slide, _SLATE_800)

        # Decorative mirrored geometry — bottom-left corner
        _add_rect(slide, 0, _SH - Inches(3.0), Inches(5.0), Inches(3.0), _INDIGO_DARK)
        _add_rect(slide, 0, _SH - Inches(2.0), Inches(3.6), Inches(2.0), _INDIGO)
        _add_rect(slide, 0, _SH - Inches(1.1), Inches(2.0), Inches(1.1), _INDIGO_MID)

        # Right accent stripe
        _add_rect(slide, _SW - Inches(0.38), 0, Inches(0.38), _SH, _INDIGO)

        # Top band
        _add_rect(slide, 0, 0, _SW - Inches(0.38), Inches(1.3), _INDIGO_DARK)

        _add_text(slide, "REPORT COMPLETE",
                  Inches(0.6), Inches(1.7), Inches(9.5), Inches(0.55),
                  size=14, bold=True, color=_SLATE_400)
        _add_text(slide, "Thank You",
                  Inches(0.6), Inches(2.2), Inches(9.5), Inches(1.3),
                  size=46, bold=True, color=_WHITE)
        _add_rect(slide, Inches(0.6), Inches(3.6), Inches(4.5), Inches(0.05), _INDIGO)
        _add_text(slide, "Gulfstream Intelligence  ·  Regulatory Intelligence Platform",
                  Inches(0.6), Inches(3.78), Inches(9.5), Inches(0.5),
                  size=12, color=_SLATE_400)
        _add_text(slide, "This presentation is confidential and for preparation purposes only.",
                  Inches(0.6), Inches(4.45), Inches(9.5), Inches(0.4),
                  size=9, color=_SLATE_600)

        buf = io.BytesIO()
        prs.save(buf)
        url = await storage_service.upload_file(
            buf.getvalue(),
            f"Report_{filename}.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        return url


export_service = ExportService()
