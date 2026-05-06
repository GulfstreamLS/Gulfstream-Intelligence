import io

import pytesseract
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader


class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

    def extract_text(self, content: bytes, file_type: str) -> str:
        """Extract text based on file type."""
        if file_type.lower() in ["pdf", "application/pdf"]:
            return self._extract_from_pdf(content)
        elif file_type.lower() in [
            "docx",
            "doc",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ]:
            return self._extract_from_docx(content)
        elif file_type.lower() in ["txt", "text/plain"]:
            return content.decode("utf-8")
        elif file_type.lower() in ["pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            return self._extract_from_pptx(content)
        elif file_type.lower() in ["png", "jpg", "jpeg", "image/png", "image/jpeg"]:
            return self._extract_from_image(content)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _extract_from_pdf(self, content: bytes) -> str:
        reader = PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        return text

    def _extract_from_docx(self, content: bytes) -> str:
        doc = Document(io.BytesIO(content))
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_from_pptx(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_from_image(self, content: bytes) -> str:
        img = Image.open(io.BytesIO(content))
        return pytesseract.image_to_string(img)

    def split_text(self, text: str) -> list[str]:
        """Split text into manageable chunks for embedding."""
        return self.text_splitter.split_text(text)


document_processor = DocumentProcessor()
